from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color palette ───
DARK = RGBColor(0x0F, 0x17, 0x2A)
SURFACE = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x1A, 0x56, 0xDB)
LIGHT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
TEAL = RGBColor(0x0E, 0xD4, 0x8F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        import lxml.etree as etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//a:srgbClr', nsmap)
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Inter'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=WHITE, alignment=PP_ALIGN.LEFT):
    """Add text box with multiple paragraphs. Each line can be a string or (text, bold, color, size) tuple."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, bold, clr, sz = line, False, color, font_size
        else:
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            clr = line[2] if len(line) > 2 else color
            sz = line[3] if len(line) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = 'Inter'
        p.alignment = alignment
        p.space_after = Pt(6)
    return txBox

def add_card(slide, left, top, width, height, title, items, accent_color=BLUE):
    card = add_shape(slide, left, top, width, height, SURFACE)
    # accent line
    add_shape(slide, left, top, Inches(0.06), height, accent_color)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.5), Inches(0.5),
                 title, 20, WHITE, True)
    add_multi_text(slide, left + Inches(0.3), top + Inches(0.6), width - Inches(0.5), height - Inches(0.8),
                   [f"  \u2713  {item}" for item in items], 14, GRAY)


# ═══════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
# gradient-like overlay
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(0x0A, 0x0F, 0x1F))
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), BLUE)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
             'UNIVERSITY COLLEGE OF ENGINEERING, ARIYALUR', 14, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
             'DormFlow', 72, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             'Dual-Hostel QR Outpass & Leave Management System', 28, LIGHT_BLUE, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             'Role-Based SPA  |  Vanilla JS + Supabase + Capacitor  |  Android APK', 16, GRAY, False, PP_ALIGN.CENTER)

# Separator
add_shape(slide, Inches(5.5), Inches(5.0), Inches(2.333), Inches(0.04), TEAL)

add_text_box(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.4),
             'Presented by MooN Software Solutions', 18, GRAY, False, PP_ALIGN.CENTER)

# Team names
add_multi_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8), [
    'Varun C  (Lead Developer)     |     Veerakumar  (Developer)     |     Umar Sheriff  (UI/UX & Testing)'
], 14, GRAY, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 2 — Team
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT_GREEN)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Meet the Team', 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             'MooN Software Solutions', 16, GRAY)

team_data = [
    ('Varun C', 'Lead Developer', 'System architecture & full-stack development.\nDesigned and built DormFlow from the ground up.', ACCENT_PURPLE),
    ('Veerakumar', 'Developer', 'Core developer — contributed to building\nthe application end-to-end.', BLUE),
    ('Umar Sheriff', 'UI/UX & Testing', 'Crafted the user experience and ensured\nquality across every screen.', TEAL),
]
for i, (name, role, desc, color) in enumerate(team_data):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(1.6), Inches(3.6), Inches(4.8), SURFACE)
    add_shape(slide, x, Inches(1.6), Inches(3.6), Inches(0.06), color)
    # Avatar circle placeholder
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), Inches(2.0), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    # Initials
    add_text_box(slide, x + Inches(1.3), Inches(2.2), Inches(1.0), Inches(0.6),
                 name[0], 32, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.2), Inches(3.2), Inches(0.4),
                 name, 24, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.6), Inches(3.2), Inches(0.3),
                 role, 14, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(4.2), Inches(3.2), Inches(1.5),
                 desc, 13, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT_RED)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Problem Statement', 36, WHITE, True)

problems = [
    'Manual paper-based leave registers are slow, error-prone, and easily lost.',
    'Wardens have no real-time visibility into who is in or out of the hostel.',
    'Gate security cannot instantly verify if a student has valid leave approval.',
    'No unified system — students, wardens, and security operate in silos.',
    '500+ hostel residents across boys & girls hostels with no digital tracking.',
]
add_multi_text(slide, Inches(1.0), Inches(1.4), Inches(11), Inches(5.5),
               [f'\u2716  {p}' for p in problems], 18, GRAY, PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════
# SLIDE 4 — Solution Overview
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Solution: DormFlow', 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             'A unified digital ecosystem for hostel leave operations', 16, GRAY)

solutions = [
    ('QR-Based Outpass System', 'Each approved leave generates a unique scannable QR pass.\nGate security validates via camera — no paper needed.'),
    ('Role-Based Access', 'Four tiers: Student, Warden, Gate Security, Admin.\nEach sees only what they need to see.'),
    ('Real-Time Status Tracking', 'Know instantly who is IN, who is OUT, and who is\npending approval — across both hostels.'),
    ('Android APK Distribution', 'Distributed as a signed APK (~5.1 MB).\nWorks on any Android 7.0+ device.'),
]
for i, (title, desc) in enumerate(solutions):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(1.6 + (i // 2) * 2.6)
    card = add_shape(slide, x, y, Inches(5.5), Inches(2.2), SURFACE)
    add_shape(slide, x, y, Inches(5.5), Inches(0.04), TEAL)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.8), Inches(0.4),
                 title, 20, WHITE, True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(4.8), Inches(1.3),
                 desc, 13, GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 5 — Student Features
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Student Features', 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             'Empowering 500+ hostel residents with digital leave management', 16, GRAY)

add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5),
         'Leave Management', [
    'Submit leave requests with date range & reason',
    'View approval status in real-time',
    'Track complete leave history',
    'Edit/cancel pending requests',
], BLUE)

add_card(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.5),
         'QR Outpass', [
    'Receive scannable QR pass on approval',
    'QR encodes leave ID + student info',
    'Present at gate for scan-in/scan-out',
], ACCENT_GREEN)

add_card(slide, Inches(6.8), Inches(4.3), Inches(5.5), Inches(2.7),
         'Profile & Dashboard', [
    'Personal profile with registration details',
    'Dashboard with leave stats & quick actions',
    'Guardian contact information on file',
], ACCENT_ORANGE)

# ═══════════════════════════════════════════════════
# SLIDE 6 — Warden Features
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT_ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Warden Features', 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             'Block-segregated approval workflow for hostel wardens', 16, GRAY)

warden_features = [
    ('Block-Segregated Queue', 'Boys & girls wardens see only their respective hostel leave requests. Zero cross-block noise.'),
    ('Approve / Reject Requests', 'One-tap approve or reject with review of student details, dates, and reason.'),
    ('Resident Directory', 'Complete view of all residents in the assigned block with room numbers and contact info.'),
    ('Leave History', 'Access full leave history for any resident in the block for audit purposes.'),
]
for i, (title, desc) in enumerate(warden_features):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(1.6 + (i // 2) * 2.6)
    card = add_shape(slide, x, y, Inches(5.5), Inches(2.2), SURFACE)
    add_shape(slide, x, y, Inches(5.5), Inches(0.04), ACCENT_ORANGE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.8), Inches(0.4),
                 title, 20, WHITE, True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(4.8), Inches(1.3),
                 desc, 13, GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 7 — Gate Security Features
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT_GREEN)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Gate Security Features', 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             'Real-time QR scanning at the hostel gate', 16, GRAY)

add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.5),
         'QR Scan', [
    'Hardware camera integration',
    'Instant pass validation',
    'Auto-detect valid/expired passes',
    '10-min re-scan cooldown',
], ACCENT_GREEN)

add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.5),
         'Gate Dashboard', [
    'Live IN/OUT status board',
    'Currently in-house residents list',
    'Scan history with timestamps',
    'Student photo verification',
], BLUE)

add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.5),
         'Security Controls', [
    'Flag invalid passes',
    'Manual registration lookup',
    'Cross-hostel gate view',
    'Audit trail for every scan',
], ACCENT_PURPLE)

# ═══════════════════════════════════════════════════
# SLIDE 8 — Admin Features
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Admin Features', 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             'Central administration & system control', 16, GRAY)

admin_data = [
    ('User Management', 'Approve/reject registrations.\nAssign roles across all 4 tiers.\nMonitor all active accounts.'),
    ('Leave Oversight', 'View ALL leave & outpass records.\nCross-hostel audit capability.\nOverride approvals when needed.'),
    ('System Dashboard', 'Real-time occupancy stats.\nActive vs pending leaves chart.\nRole-wise user distribution.'),
    ('Registration Gate', 'Single source of truth for access.\nValidate new users across roles.\nBlock unauthorized accounts.'),
]
for i, (title, desc) in enumerate(admin_data):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(1.6 + (i // 2) * 2.6)
    card = add_shape(slide, x, y, Inches(5.5), Inches(2.2), SURFACE)
    add_shape(slide, x, y, Inches(5.5), Inches(0.04), ACCENT_PURPLE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.8), Inches(0.4),
                 title, 20, WHITE, True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(4.8), Inches(1.3),
                 desc, 13, GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 9 — Tech Stack
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Technology Stack', 36, WHITE, True)

stacks = [
    ('Frontend', [
        'Vanilla JavaScript (ES Modules)',
        'Vite Build Tool',
        'Chart.js (Dashboards)',
        'html5-qrcode (QR Scanner)',
        'qrcode (QR Generator)',
        'Fontsource Inter + Material Icons',
    ]),
    ('Backend & DB', [
        'Supabase (BaaS)',
        'PostgreSQL with RLS Policies',
        'Supabase Auth (Email + Password)',
        'Row-Level Security (RLS)',
        'Automatic profile triggers',
    ]),
    ('Mobile', [
        'Apache Capacitor v8',
        'Android APK (v2 signed)',
        'UsesCleartextTraffic enabled',
        'Min SDK: Android 7.0',
        'APK Size: ~5.1 MB',
    ]),
    ('DevOps & Hosting', [
        'GitHub (Source Control)',
        'GitHub Pages (Landing Page)',
        'Supabase Cloud (Hosted DB)',
        'Supabase Studio (Admin UI)',
        'SHA-256 code signing',
    ]),
]
for i, (title, items) in enumerate(stacks):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(1.2 + (i // 2) * 3.0)
    card = add_shape(slide, x, y, Inches(5.5), Inches(2.6), SURFACE)
    add_shape(slide, x, y, Inches(5.5), Inches(0.04), BLUE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.8), Inches(0.4),
                 title, 20, WHITE, True)
    add_multi_text(slide, x + Inches(0.3), y + Inches(0.65), Inches(4.8), Inches(1.8),
                   [f'  \u25B8  {item}' for item in items], 13, GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 10 — Architecture
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'System Architecture', 36, WHITE, True)

# Architecture diagram using shapes
layers = [
    ('CLIENT LAYER', 'Vanilla JS SPA\nVite Bundle\nAndroid WebView', Inches(0.8), BLUE),
    ('AUTH LAYER', 'Supabase Auth\nJWT Tokens\nRole Verification', Inches(2.8), ACCENT_PURPLE),
    ('API LAYER', 'Supabase REST API\nRealtime Subscriptions\nStorage API', Inches(4.8), ACCENT_ORANGE),
    ('DATABASE LAYER', 'PostgreSQL\nRLS Policies\nTriggers & Funcs', Inches(6.8), ACCENT_GREEN),
    ('STORAGE LAYER', 'Supabase Storage\nAvatars / Photos\nPublic Assets', Inches(8.8), TEAL),
]

for label, desc, x, color in layers:
    card = add_shape(slide, x, Inches(1.4), Inches(3.6), Inches(2.8), SURFACE)
    add_shape(slide, x, Inches(1.4), Inches(3.6), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(3.2), Inches(0.3),
                 label, 13, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.1), Inches(3.2), Inches(1.8),
                 desc, 14, GRAY, False, PP_ALIGN.CENTER)

# Arrows (using text)
for i in range(4):
    x = Inches(1.0 + i * 2.0)
    add_text_box(slide, Inches(4.35 + i * 2.0), Inches(2.5), Inches(0.5), Inches(0.4),
                 '\u25B6', 24, GRAY, False, PP_ALIGN.CENTER)

# Data flow description
add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.5), SURFACE)
add_text_box(slide, Inches(1.1), Inches(4.75), Inches(11), Inches(0.3),
             'Data Flow', 18, WHITE, True)

flow_lines = [
    '1.  User opens app \u2192 SPA loads in Android WebView \u2192 checks for existing session',
    '2.  Login/Register \u2192 Supabase Auth validates credentials \u2192 returns JWT with role in app_metadata',
    '3.  Every API request \u2192 JWT sent as Bearer token \u2192 PostgreSQL RLS enforces row-level permissions',
    '4.  Leave request \u2192 stored in leaves table \u2192 triggers outpass generation \u2192 QR code created on approval',
    '5.  Gate scan \u2192 camera reads QR \u2192 out_time/in_time logged \u2192 10-min cooldown enforced server-side',
]
add_multi_text(slide, Inches(1.1), Inches(5.15), Inches(11), Inches(1.8),
               flow_lines, 12, GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 11 — Database Schema
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Database Schema & Security', 36, WHITE, True)

tables = [
    ('profiles', [
        'id (UUID, PK)  |  email  |  name',
        'role (student/boys_warden/girls_warden/security/admin)',
        'gender  |  hostel_type  |  department  |  year',
        'room_number  |  block_name  |  phone',
        'guardian_name  |  guardian_phone',
        'active_status (IN/OUT)  |  is_approved',
    ]),
    ('leaves', [
        'id (UUID, PK)  |  student_id (FK)',
        'start_date  |  end_date  |  reason',
        'approval_status (Pending/Approved/Rejected)',
        'approved_by (FK)  |  applied_at  |  updated_at',
    ]),
    ('outpasses', [
        'id (UUID, PK)  |  leave_id (FK, unique)',
        'student_id (FK)  |  qr_code',
        'out_time  |  in_time',
        'scanned_by (FK)  |  created_at',
    ]),
]

for i, (table, cols) in enumerate(tables):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(1.4), Inches(3.6), Inches(4.5), SURFACE)
    add_shape(slide, x, Inches(1.4), Inches(3.6), Inches(0.06), BLUE)
    add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(3.2), Inches(0.3),
                 table, 18, WHITE, True, PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.2), Inches(2.1), Inches(3.2), Inches(3.5),
                   cols, 12, GRAY)

# Security note
add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8), SURFACE)
add_shape(slide, Inches(0.8), Inches(6.2), Inches(0.06), Inches(0.8), ACCENT_GREEN)
add_text_box(slide, Inches(1.2), Inches(6.3), Inches(11), Inches(0.6),
             'Security: Row-Level Security (RLS) policies on every table. JWT role from app_metadata controls all CRUD operations. '
             'Triggers auto-create profiles on signup. Unique constraint on outpasses.leave_id ensures 1:1 mapping.', 12, GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 12 — Security
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT_RED)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Security Architecture', 36, WHITE, True)

sec_items = [
    ('Row-Level Security (RLS)', 'PostgreSQL policies on every table. Students see only their own data.\nWardens see only their block. Admins see everything. Enforced at DB level.'),
    ('JWT Authentication', 'Supabase Auth with email/password. Role stored in app_metadata\nwithin JWT. Every API request verified server-side.'),
    ('Admin Approval Gate', 'New registrations set is_approved=false. Admin must manually\napprove before login is allowed. Prevents unauthorized access.'),
    ('QR Code Validation', 'Outpass QR encodes leave UUID + student ID. Gate scan validates\nagainst DB in real-time. Expired/invalid passes flagged immediately.'),
    ('10-Minute Scan Cooldown', 'After OUT scan, gate cannot re-scan for IN within 10 minutes.\nPrevents duplicate scans. Enforced server-side in store.js.'),
    ('APK Signing', 'Signed with SHA-256 RSA 2048-bit keystore.\nv2 APK Signature Scheme. Secure distribution.'),
]
for i, (title, desc) in enumerate(sec_items):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(1.3 + (i // 2) * 2.0)
    card = add_shape(slide, x, y, Inches(5.5), Inches(1.7), SURFACE)
    add_shape(slide, x, y, Inches(0.06), Inches(1.7), ACCENT_RED)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(4.8), Inches(0.3),
                 title, 16, WHITE, True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.5), Inches(4.8), Inches(1.1),
                 desc, 12, GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 13 — Workflow / User Flow
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'User Workflow', 36, WHITE, True)

steps = [
    ('1', 'Student\nSubmits Leave', BLUE),
    ('2', 'Warden\nApproves', ACCENT_ORANGE),
    ('3', 'QR Pass\nGenerated', ACCENT_PURPLE),
    ('4', 'Gate Scan\nOUT', ACCENT_GREEN),
    ('5', 'Gate Scan\nIN', TEAL),
    ('6', 'Auto\nCheck-In', ACCENT_PURPLE),
]
for i, (num, label, color) in enumerate(steps):
    x = Inches(0.8 + i * 2.1)
    # Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), Inches(1.6), Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.35), Inches(1.75), Inches(1.2), Inches(0.9),
                 num, 36, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.0), Inches(1.9), Inches(0.6),
                 label, 14, WHITE, True, PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(1.5), Inches(2.0), Inches(0.6), Inches(0.4),
                     '\u25B6', 20, GRAY, False, PP_ALIGN.CENTER)

# Detail cards
details = [
    'Student logs in, fills leave form with\ndates & reason. Request enters\nPending queue.',
    'Warden sees request in block-specific\nqueue. Reviews & taps Approve\nor Reject.',
    'System generates QR code\nencoding leave ID + student\nID. Available in app.',
    'Gate security scans QR at exit.\nSystem logs out_time.\nCooldown starts.',
    'Student returns, gate scans QR.\nSystem logs in_time.\nAuto-updates status.',
    'Active status set to IN.\nLeave marked complete.\nDashboard updated.',
]
for i, detail in enumerate(details):
    x = Inches(0.5 + i * 2.1)
    y = Inches(3.8)
    card = add_shape(slide, x, y, Inches(1.9), Inches(3.2), SURFACE)
    add_shape(slide, x, y, Inches(1.9), Inches(0.04), GRAY)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(1.7), Inches(2.9),
                 detail, 11, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 14 — Live Demo / Deployment
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT_GREEN)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Deployment & Availability', 36, WHITE, True)

deploy_items = [
    ('Android APK', 'Signed release APK (v3.0.0)\nSize: ~5.1 MB\nAndroid 7.0+ required\nv2 APK Signature Scheme', ACCENT_GREEN),
    ('GitHub Repository', 'Source code hosted at:\ngithub.com/VARUN122134/dormflow\n98 files, 10K+ lines of code\nFull commit history', BLUE),
    ('GitHub Pages Landing', 'Download page at:\nvarun122134.github.io/dormflow/\nFeature showcase + install guide\nDirect APK download', ACCENT_ORANGE),
    ('Supabase Backend', 'Cloud-hosted PostgreSQL\nRow-Level Security enabled\nSupabase Auth + Storage\n50K MAU free tier capacity', ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(deploy_items):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(1.4 + (i // 2) * 2.8)
    card = add_shape(slide, x, y, Inches(5.5), Inches(2.4), SURFACE)
    add_shape(slide, x, y, Inches(5.5), Inches(0.04), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.8), Inches(0.3),
                 title, 22, WHITE, True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.65), Inches(4.8), Inches(1.5),
                 desc, 13, GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 15 — Future Scope
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Future Enhancements', 36, WHITE, True)

future = [
    'Real-time WebSocket notifications for leave status changes',
    'Multi-campus support with role-based campus assignment',
    'Push notifications via Firebase Cloud Messaging',
    'Forgot password & email verification flow',
    'Post Announcements module for wardens',
    'AI-based leave approval suggestions using historical data',
    'iOS version via Capacitor (minimal code changes)',
    'Biometric authentication (fingerprint / face unlock)',
    'Offline-first architecture with local data sync',
    'QR gate hardware integration (Raspberry Pi scanners)',
]
add_multi_text(slide, Inches(1.0), Inches(1.2), Inches(11), Inches(5.5),
               [f'  \u25B8  {item}' for item in future], 17, GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 16 — Stats / Impact
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             'Project Impact', 36, WHITE, True)

stats = [
    ('500+', 'Students Served\n(Boys + Girls Hostels)', BLUE),
    ('4', 'User Roles\n(Student/Warden/Security/Admin)', ACCENT_ORANGE),
    ('10K+', 'Lines of Code\n(Vanilla JS Architecture)', ACCENT_PURPLE),
    ('5.1 MB', 'APK Size\n(Minimal Footprint)', ACCENT_GREEN),
    ('Zero', 'Paper Usage\n(Fully Digital Workflow)', TEAL),
    ('50K', 'Monthly Active Users\n(Supabase Free Tier Capacity)', ACCENT_RED),
]
for i, (num, label, color) in enumerate(stats):
    x = Inches(0.8 + (i % 3) * 4.0)
    y = Inches(1.5 + (i // 3) * 2.6)
    card = add_shape(slide, x, y, Inches(3.6), Inches(2.2), SURFACE)
    add_shape(slide, x, y, Inches(3.6), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.2), Inches(0.7),
                 num, 42, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(3.2), Inches(0.8),
                 label, 14, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 17 — Thank You / Q&A
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), TEAL)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             'Thank You', 72, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.6),
             'Questions & Discussion', 28, LIGHT_BLUE, False, PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.8), Inches(2.333), Inches(0.04), TEAL)

add_multi_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(2.0), [
    'University College of Engineering, Ariyalur',
    'MooN Software Solutions  |  DormFlow v3.0.0',
    'github.com/VARUN122134/dormflow',
], 16, GRAY, PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'DormFlow_Presentation.pptx')
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
